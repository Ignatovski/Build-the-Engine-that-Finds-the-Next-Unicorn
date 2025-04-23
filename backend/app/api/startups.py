from fastapi import APIRouter, HTTPException, Query, File, Form, UploadFile
from typing import Dict, Optional, List, Any
from pydantic import BaseModel
from PIL import Image
from PyPDF2 import PdfReader
from io import BytesIO
import aiohttp
import os
import json
import tempfile
import openai
import ascii_magic

from ..schemas.startup import StartupAnalysisRequest, StartupCreate
from ..services.analysis import StartupAnalyzer
from dotenv import load_dotenv

# Load environment variables from .env file
load_dotenv()

# Configure API keys
openai_api_key = os.getenv('OPENAI_API_KEY')
serp_api_key = os.getenv('SERP_API_KEY')
tavily_api_key = os.getenv('TAVILY_API_KEY')

# Check if Tavily API key is set
if not os.getenv('TAVILY_API_KEY'):
    raise HTTPException(
        status_code=500,
        detail="Tavily API key not set. Please set TAVILY_API_KEY environment variable."
    )

# Configure OpenAI
if not os.getenv('OPENAI_API_KEY'):
    raise HTTPException(
        status_code=500,
        detail="OpenAI API key not set. Please set OPENAI_API_KEY environment variable."
    )

# Set OpenAI API key
openai.api_key = os.getenv('OPENAI_API_KEY')

# Initialize router and analyzer
router = APIRouter()
analyzer = StartupAnalyzer()

async def process_uploaded_file(file: UploadFile) -> str:
    """Process uploaded PDF or PPTX file and extract text content"""
    print(f"[DEBUG] Processing uploaded file: {file.filename}")
    content = ""
    
    try:
        # Read file content
        file_content = await file.read()
        
        if file.filename.lower().endswith('.pdf'):
            print("[DEBUG] Processing PDF file")
            # Process PDF
            pdf_reader = PdfReader(BytesIO(file_content))
            for page in pdf_reader.pages:
                content += page.extract_text() + "\n"
                
        elif file.filename.lower().endswith('.pptx'):
            print("[DEBUG] Processing PPTX file")
            # For PPTX we'll need python-pptx
            from pptx import Presentation
            prs = Presentation(BytesIO(file_content))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
        else:
            raise ValueError(f"Unsupported file type: {file.filename}")
            
        print(f"[DEBUG] Extracted {len(content)} characters of text")
        return content
        
    except Exception as e:
        print(f"[ERROR] Failed to process file: {str(e)}")
        import traceback
        print(traceback.format_exc())
        raise HTTPException(
            status_code=400,
            detail=f"Failed to process file: {str(e)}"
        )

async def display_image_in_terminal(session: aiohttp.ClientSession, image_url: str) -> None:
    print(f"\n[DEBUG] Attempting to display image from URL: {image_url}")
    try:
        # Download the image
        print("[DEBUG] Downloading image...")
        async with session.get(image_url) as response:
            print(f"[DEBUG] Download status: {response.status}")
            if response.status != 200:
                print(f"Failed to download image: {response.status}")
                return
            
            image_data = await response.read()
            print(f"[DEBUG] Downloaded {len(image_data)} bytes")
            
            # Create a temporary file to save the image
            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
                print(f"[DEBUG] Saving to temp file: {tmp_file.name}")
                tmp_file.write(image_data)
                tmp_file.flush()
                
                try:
                    print("[DEBUG] Converting to ASCII art...")
                    # Convert to ASCII art with a reasonable size
                    output = ascii_magic.from_image_file(
                        tmp_file.name,
                        columns=60,  # Adjust width as needed
                        mode=ascii_magic.Modes.ASCII
                    )
                    
                    # Print with a border
                    width = 62  # columns + 2 for borders
                    print("\n=== Image Preview ===")
                    print("─" * width)
                    print(output)
                    print("─" * width)
                    print("=== End Preview ===")
                    
                except Exception as e:
                    print(f"[DEBUG] Failed to convert image to ASCII: {str(e)}")
                    import traceback
                    print(traceback.format_exc())
                
                # Clean up
                print("[DEBUG] Cleaning up temp file...")
                os.unlink(tmp_file.name)
                
    except Exception as e:
        print(f"[DEBUG] Error displaying image: {str(e)}")
        import traceback
        print(traceback.format_exc())

async def search_web_data(startup_name: str, industry: Optional[str] = None):
    """Search for company information using Tavily API"""
    # Get API key from environment
    api_key = os.getenv('TAVILY_API_KEY')
    if not api_key:
        raise ValueError("TAVILY_API_KEY environment variable not set")
    
    # API endpoint and headers
    url = "https://api.tavily.com/search"
    headers = {
        'Content-Type': 'application/json',
        'Authorization': f'Bearer {api_key}'
    }
    
    # Prepare search query
    search_data = {
        "query": startup_name,
        "search_depth": "advanced",
        "max_results": 10,
        "include_answer": True
    }
    
    website_url = None
    try:
        async with aiohttp.ClientSession() as session:
            async with session.post(url, json=search_data, headers=headers) as response:
                if response.status == 200:
                    response_data = await response.json()
                    
                    if not isinstance(response_data, dict):
                        raise ValueError("Invalid response format from Tavily API")
                    
                    # Extract website URL from results
                    if 'results' in response_data:
                        for result in response_data['results']:
                            if isinstance(result, dict):
                                url_str = result.get('url', '')
                                if url_str and startup_name.lower() in url_str.lower():
                                    website_url = url_str
                                    break
                    
                    # Process search results
                    search_results = []
                    relevant_info = []
                    
                    if 'results' in response_data:
                        for result in response_data['results']:
                            if isinstance(result, dict):
                                title = result.get('title', '')
                                content = result.get('content', '')
                                if title and content:
                                    search_results.append({
                                        'title': title,
                                        'content': content
                                    })
                                    relevant_info.append(f"Title: {title}\nContent: {content}\n")
                    
                    search_insights = "\n".join(relevant_info) if relevant_info else "No relevant information found."
                    
                    return {
                        'data': response_data,
                        'website_url': website_url,
                        'logo_url': None,
                        'search_insights': search_insights,
                        'search_results': search_results
                    }
                else:
                    error_text = await response.text()
                    raise HTTPException(
                        status_code=response.status,
                        detail=f"Tavily API error: {error_text}"
                    )
    except aiohttp.ClientError as e:
        raise HTTPException(status_code=502, detail=f"Network error: {str(e)}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Unexpected error: {str(e)}")
    # If we didn't get any successful response
    return {
        "data": {},
        "website_url": None,
        "logo_url": None,
        "search_insights": "No information found.",
        "search_results": []
    }

@router.post("/analyze")
async def analyze_startup(
    file: Optional[UploadFile] = File(None),
    data: str = Form(...)
):
    """Analyze a startup's potential using Tavily search and OpenAI"""
    # Parse the JSON data
    startup_data = json.loads(data)
    startup = StartupAnalysisRequest(**startup_data)
    
    print(f"Analyzing startup: {startup.name}")
    print(f"File uploaded: {file.filename if file else 'No file'}")
    
    try:
        # Process uploaded file if present
        file_content = ""
        if file:
            try:
                file_content = await process_uploaded_file(file)
                print(f"Successfully processed file: {file.filename}")
            except Exception as e:
                print(f"Error processing file: {str(e)}")
                # Continue with analysis even if file processing fails
        
        # Get web data about the startup
        search_response = await search_web_data(startup.name, startup.industry)
        
        # Extract web data and website URL from response
        web_data = search_response.get('data', {})
        website_url = search_response.get('website_url')
        logo_url = search_response.get('logo_url')
        search_insights = search_response.get('search_insights', 'No relevant information found.')
        
        # Prepare prompt for OpenAI analysis
        prompt = (
            f"Analyze the following startup:\n\n"
            f"Name: {startup.name}\n"
            f"Industry: {startup.industry}\n"
            f"Description: {startup.description}\n\n"
            f"Web Research Results:\n{search_insights}\n\n"
            f"Uploaded Document Content:\n{file_content}\n\n"
            f"Based on the available information, provide a comprehensive analysis covering:\n"
            f"1. Market Score (0-10): Market size, growth potential, competition\n"
            f"2. Team Score (0-10): Leadership, experience, track record\n"
            f"3. Technology Score (0-10): Innovation, technical advantage, IP\n"
            f"4. Traction Score (0-10): Growth, user base, revenue\n"
            f"5. Overall Score (0-10): Weighted average\n"
            f"6. Success Probability (0-1): Likelihood of significant growth\n"
            f"7. Unicorn Probability (0-1): Chance of reaching $1B valuation\n"
            f"8. BCG Matrix Parameters:\n"
            f"   - Market Growth (0-1): Industry growth rate, YoY growth\n"
            f"   - Market Potential (0-1): Relative market share potential\n"
            f"9. Key Strengths (3-5 bullet points)\n"
            f"10. Major Risks (2-3 bullet points)\n"
            f"11. Strategic Recommendations (2-3 bullet points)\n\n"
            f"Format your response as structured JSON with keys: market_score, "
            f"team_score, technology_score, traction_score, overall_score, "
            f"success_probability, unicorn_probability, market_growth, market_potential, "
            f"strengths, risks, recommendations.\n\n"
            f"Note: market_growth and market_potential should be between 0 and 1, "
            f"representing the startup's position on the BCG Matrix."
        )
        try:
            client = openai.AsyncOpenAI()
            response = await client.chat.completions.create(
                model="gpt-4",
                messages=[
                    {"role": "system", "content": "You are a startup analyst. Respond ONLY in valid JSON with all numeric values as numbers, not strings."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.7
            )
            
            # Get the response content
            content = response.choices[0].message.content
            
            # Parse the JSON response
            analysis = json.loads(content)
            
            # Add website and logo URLs
            analysis['website_url'] = website_url
            analysis['logo_url'] = logo_url
            
            return {
                'analysis': analysis
            }
            
        except json.JSONDecodeError as e:
            raise HTTPException(status_code=500, detail=f"Failed to parse OpenAI response: {str(e)}")

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error during analysis: {str(e)}")

async def create_startup(startup: StartupCreate):
    """Create a new startup entry and analyze it using AI."""
    try:
        # Analyze startup using OpenAI
        prompt = (
            f"You are analyzing {startup.name}. Use all available information to provide a startup analysis.\n\n"
            f"Company Details:\n"
            f"- Name: {startup.name}\n"
            f"- Description: {startup.description or 'Not provided'}\n"
            f"- Industry: {startup.industry or 'Not provided'}\n"
            f"- Funding Stage: {startup.funding_stage or 'Not provided'}\n"
            f"- Website: {startup.website or 'Not provided'}\n\n"
            f"Web Research Results:\n{search_insights}\n\n"
            f"Uploaded Document Content:\n{file_content}\n\n"
            f"Based on the available information, provide a comprehensive analysis covering:\n"
            f"1. Market Score (0-10): Market size, growth potential, competition\n"
            f"2. Team Score (0-10): Leadership, experience, track record\n"
            f"3. Technology Score (0-10): Innovation, technical advantage, IP\n"
            f"4. Traction Score (0-10): Growth, user base, revenue\n"
            f"5. Overall Score (0-10): Weighted average\n"
            f"6. Success Probability (0-1): Likelihood of significant growth\n"
            f"7. Unicorn Probability (0-1): Chance of reaching $1B valuation\n"
            f"8. BCG Matrix Parameters:\n"
            f"   - Market Growth (0-1): Industry growth rate, YoY growth\n"
            f"   - Market Potential (0-1): Relative market share potential\n"
            f"9. Key Strengths (3-5 bullet points)\n"
            f"10. Major Risks (2-3 bullet points)\n"
            f"11. Strategic Recommendations (2-3 bullet points)\n\n"
            f"Format your response as structured JSON with keys: market_score, "
            f"team_score, technology_score, traction_score, overall_score, "
            f"success_probability, unicorn_probability, market_growth, market_potential, "
            f"strengths, risks, recommendations.\n\n"
            f"Note: market_growth and market_potential should be between 0 and 1, "
            f"representing the startup's position on the BCG Matrix."
        )
        try:
            response = await openai.ChatCompletion.acreate(
                model="gpt-4",
                messages=[
                    {"role": "system", "content": "You are a startup analyst. Respond ONLY in valid JSON with all numeric values as numbers, not strings."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.7
            )

            # Get the response content
            content = response.choices[0].message.content
            print("OpenAI Response:", content)
            
            # Parse the JSON response
            analysis = json.loads(content)
            
            # Handle website URL
            analysis['website_url'] = None  # Initialize as None
            
            # Add website URL
            analysis['website_url'] = startup.website or website_url or None

            # Add logo URL to analysis
            analysis['logo_url'] = logo_url
            
            # Validate URL format if present
            if analysis['website_url']:
                try:
                    if not analysis['website_url'].startswith(('http://', 'https://')):
                        analysis['website_url'] = 'https://' + analysis['website_url']
                except:
                    analysis['website_url'] = None
            
            # Ensure all required fields are present and numeric
            required_numeric_fields = [
                'market_score', 'team_score', 'technology_score', 'traction_score',
                'overall_score', 'success_probability', 'unicorn_probability',
                'market_growth', 'market_potential'
            ]
            
            for field in required_numeric_fields:
                if field not in analysis:
                    analysis[field] = 0.5  # Default value
                elif not isinstance(analysis[field], (int, float)):
                    try:
                        analysis[field] = float(analysis[field])
                    except (ValueError, TypeError):
                        analysis[field] = 0.5  # Default value
            
            # Ensure array fields are present
            array_fields = ['strengths', 'risks', 'recommendations']
            for field in array_fields:
                if field not in analysis or not isinstance(analysis[field], list):
                    analysis[field] = []
                # Fallback: return the raw content if parsing fails
                analysis = {"raw": raw_content}

            return {"startup": startup.dict(), "analysis": analysis}

        except openai.error.OpenAIError as e:
            raise HTTPException(status_code=503, detail=f"OpenAI API error: {str(e)}")
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Analysis error: {str(e)}")
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Invalid request: {str(e)}")

@router.get("/startups/", response_model=List[Dict])
async def get_startups():
    """Get all startups with analysis (real version: returns empty until DB integration)"""
    # TODO: Integrate with persistent storage to return real analyzed startups
    return []

@router.get("/startups/{startup_id}", response_model=dict)
async def get_startup(startup_id: int):
    """Get a specific startup by ID (example data)"""
    mock_profiles = {
        1: {
            "name": "AI Health Solutions",
            "description": "A pioneering healthcare startup using artificial intelligence for early disease detection.",
            "industry": "Healthcare",
            "tech_stack": ["Python", "TensorFlow", "AWS"],
            "valuation": "$2.5B",
            "team": "50+ AI researchers and healthcare professionals",
            "metrics": "500k users, 98% accuracy",
            "growth": "200% YoY user growth"
        },
        2: {
            "name": "Green Energy Tech",
            "description": "Developing next-generation solar panels with 40% higher efficiency.",
            "industry": "Clean Tech",
            "tech_stack": ["Material Science", "IoT", "Cloud Computing"],
            "valuation": "$1.8B",
            "team": "100+ engineers and scientists",
            "metrics": "300MW installed capacity",
            "growth": "150% YoY revenue growth"
        }
    }
    
    if startup_id in mock_profiles:
        return {"id": startup_id, "profile": mock_profiles[startup_id]}
    else:
        return {"id": startup_id, "profile": mock_profiles[1]}

@router.get("/startups/search/ai", response_model=List[Dict])
async def search_startups(
    query: str = Query(..., description="Search query"),
    industry: Optional[str] = Query(None, description="Industry filter")
):
    """Search startups using Tavily API (real data)"""
    if not tavily_api_key:
        return [{"error": "Tavily API key not configured"}]
    try:
        async with aiohttp.ClientSession() as session:
            headers = {"Authorization": f"Bearer {tavily_api_key}"}
            q = query if not industry else f"{query} {industry} startup"
            params = {
                "query": q,
                "search_depth": "advanced",
                "include_domains": [
                    "techcrunch.com", "crunchbase.com", "linkedin.com", "bloomberg.com", "reuters.com", "forbes.com"
                ]
            }
            print(f"Making request to Tavily API with query: {q}")
            async with session.get(
                "https://api.tavily.com/search",
                headers={"api-key": tavily_api_key},
                params=params
            ) as response:
                if response.status == 200:
                    tavily_data = await response.json()
                    print(f"Tavily API response: {tavily_data}")
                    # Structure the response for the frontend
                    results = []
                    if "results" in tavily_data:
                        for item in tavily_data["results"]:
                            result = {
                                "title": item.get("title"),
                                "url": item.get("url"),
                                "snippet": item.get("content"),  # Tavily uses 'content'
                                "source": item.get("source")
                            }
                            if all(result.values()):
                                results.append(result)
                    if not results:
                        print("No valid results found in Tavily response")
                    return [{"synthetic": False, "results": results}]
                else:
                    error_text = await response.text()
                    print(f"Tavily API error: {response.status}, {error_text}")
                    return [{"error": f"Tavily API error: {response.status}"}]
    except Exception as e:
        import traceback
        print(f"Error in Tavily search: {str(e)}")
        print(f"Traceback: {traceback.format_exc()}")
        return [{"error": str(e)}]
